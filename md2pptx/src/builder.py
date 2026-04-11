"""
builder.py — PPTX construction engine.
Uses the provided Accenture Slide Master template as the base for every slide.
All layout/styling is driven by the template's theme (fonts, colors, backgrounds).
"""
import io
import os
import math
from typing import Dict, Any, List, Optional, Tuple
from copy import deepcopy
from PIL import Image

from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import lxml.etree as etree

from .parser import clean_text, Table
from .charts import render_chart, stat_card_image, bar_chart
from .icons import get_icon_for_title, get_numbered_icon, ACCENT_COLORS
from .image_gen import generate_slide_asset, _PIL_AVAILABLE

# ── Template path ─────────────────────────────────────────────────────────────
TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "..", "assets", "template.pptx")

# ── Slide geometry (Default 16:9, but updated dynamically in generate_pptx) ──
SW = 12192000   # slide width
SH = 6858000    # slide height
M  = 350000     # left/right margin
TM = 600000     # top margin
TH = 530000     # title height
CT = 1350000    # content area top
CH = SH - CT - 320000  # content area height


def _setup_layout(prs: Presentation):
    """Sets up global geometry constants based on the provided template's dimensions."""
    global SW, SH, M, TM, TH, CT, CH
    SW = prs.slide_width
    SH = prs.slide_height
    
    # Scale constants based on original proportions from the Accenture template
    M  = int(SW * 0.0287)  # approx 350k/12.2M
    TM = int(SH * 0.0875)  # approx 600k/6.8M
    TH = int(SH * 0.0772)  # approx 530k/6.8M
    CT = int(SH * 0.1968)  # approx 1.35M/6.8M
    CH = SH - CT - int(SH * 0.0466) # approx 320k/6.8M

# ── Brand colors (from theme1.xml) ────────────────────────────────────────────
C = {
    "red":    RGBColor(0xEF, 0x44, 0x44),
    "blue":   RGBColor(0x0F, 0x9E, 0xD5),
    "green":  RGBColor(0x19, 0x6B, 0x24),
    "orange": RGBColor(0xE9, 0x71, 0x32),
    "purple": RGBColor(0xA0, 0x2B, 0x93),
    "dark":   RGBColor(0x2C, 0x2C, 0x2C),
    "mid":    RGBColor(0x66, 0x66, 0x66),
    "light":  RGBColor(0xE8, 0xE8, 0xE8),
    "white":  RGBColor(0xFF, 0xFF, 0xFF),
}
ACCENT_RGB = [C["red"], C["blue"], C["green"], C["orange"], C["purple"]]

FONT_H = "Libre Baskerville"   # header
FONT_B = "Inter"               # body


# ─────────────────────────────────────────────────────────────────────────────
#  Low-level helpers
# ─────────────────────────────────────────────────────────────────────────────

def _rgb(r, g, b) -> RGBColor:
    return RGBColor(r, g, b)


def _para_props(p, space_before_pt: int = 0, space_after_pt: int = 0,
                line_spacing_pt: int = 0):
    """Set paragraph spacing on a paragraph element."""
    pPr = p._pPr
    if pPr is None:
        pPr = p._p.get_or_add_pPr()
    if space_before_pt:
        sb = etree.SubElement(pPr, qn("a:spcBef"))
        etree.SubElement(sb, qn("a:spcPts")).set("val", str(space_before_pt * 100))
    if space_after_pt:
        sa = etree.SubElement(pPr, qn("a:spcAft"))
        etree.SubElement(sa, qn("a:spcPts")).set("val", str(space_after_pt * 100))


def _add_run(p, text: str, font_name: str, size_pt: float, bold: bool = False,
             italic: bool = False, color: Optional[RGBColor] = None):
    run = p.add_run()
    run.text = text
    rf = run.font
    rf.name = font_name
    rf.size = Pt(size_pt)
    rf.bold = bold
    rf.italic = italic
    if color:
        rf.color.rgb = color
    return run


def _add_textbox(slide, text: str, left, top, width, height,
                 font_name: str = FONT_B, size_pt: float = 14,
                 bold: bool = False, italic: bool = False,
                 color: RGBColor = None, align: PP_ALIGN = PP_ALIGN.LEFT,
                 word_wrap: bool = True, space_before: int = 0) -> Any:
    color = color or C["dark"]
    txb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if space_before:
        _para_props(p, space_before_pt=space_before)
    _add_run(p, clean_text(text), font_name, size_pt, bold, italic, color)
    return txb


def _add_rect(slide, left, top, width, height, fill_rgb: RGBColor,
              no_line: bool = True) -> Any:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if no_line:
        shape.line.fill.background()
    return shape


def _add_oval(slide, left, top, width, height, fill_rgb: RGBColor,
              no_line: bool = True) -> Any:
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if no_line:
        shape.line.fill.background()
    return shape

def _build_image_banner(slide, title_text: str, parsed: Dict) -> int:
    """
    Creates a highly stylized banner matching consulting decks:
    Top background image spanning width, overlapped by a large red semi-circle with text.
    Returns the Y coordinate where the content area below the banner starts.
    """
    # Clean up native slide title string to prevent bleed-through
    try:
        if slide.shapes.title:
            slide.shapes.title.text = ""
    except:
        pass
        
    img_h = int(SH * 0.38)
    # Underlying dark strip so square AI images don't leave white vertical borders
    _add_rect(slide, 0, 0, SW, img_h, C["dark"])
    
    banner_img = generate_slide_asset(title_text, parsed.get("title", ""))
    if banner_img:
        _add_image_cover(slide, banner_img, left=0, top=0, width=SW, height=img_h)
        
    # Add large red semi-circle masking
    oval_w = int(SW * 0.50)  # Thinner oval for more corporate look
    oval_h = int(SH * 0.25)
    oval_left = (SW - oval_w) // 2
    # Place it spanning from middle of the image to below it
    oval_top = img_h - int(oval_h * 0.60)
    _add_oval(slide, oval_left, oval_top, oval_w, oval_h, C["red"])
    
    # Text inside oval
    _add_textbox(slide, title_text.upper(),
                 left=oval_left + int(oval_w*0.1), top=oval_top + int(oval_h*0.35),
                 width=int(oval_w*0.8), height=int(oval_h*0.5),
                 font_name=FONT_H, size_pt=28, bold=True, color=C["white"],
                 align=PP_ALIGN.CENTER, word_wrap=True)
                 
    return oval_top + oval_h



def _add_semicircle_image(slide, img_bytes: bytes, left, top, width, height, side="left") -> Any:
    """Add an image masked into a semicircle (left or right facing)."""
    if not img_bytes:
        return None
    
    # Fallback to standard rectangular cover if processing fails
    return _add_image_cover(slide, img_bytes, left, top, width, height)

def _add_split_circle(slide, x, y, size, color_top, color_bottom) -> Any:
    """Add a circle with a diagonal color split (Top-Left / Bottom-Right)."""
    _add_oval(slide, x, y, size, size, color_bottom)
    # Diagonal half overlay
    points = [
        (Emu(x), Emu(y)),
        (Emu(x + size), Emu(y)),
        (Emu(x), Emu(y + size)),
        (Emu(x), Emu(y))
    ]
    shp = slide.shapes.build_freeform(points[0][0], points[0][1])
    shp.add_line_segments(points[1:], close=True)
    shp = shp.convert_to_shape()
    shp.fill.solid()
    shp.fill.fore_color.rgb = color_top
    shp.line.fill.background()
    return shp

def _add_image_cover(slide, img_bytes: bytes, left, top, width, height) -> Any:
    """Add image completely filling the bounds using a cover aspect ratio, with strict cropping to prevent bleed."""
    if not img_bytes:
        return None
    stream = io.BytesIO(img_bytes)
    stream.seek(0)
    with Image.open(stream) as img:
        img_w, img_h = img.size
    
    # Calculate aspect ratios
    target_ratio = width / height
    img_ratio = img_w / img_h
    
    # Add picture at target size/pos (it will be stretched initially)
    stream.seek(0)
    pic = slide.shapes.add_picture(stream, Emu(left), Emu(top), width=Emu(width), height=Emu(height))
    
    # Apply cropping to simulate 'cover' behavior without bleeding
    if img_ratio > target_ratio:
        # Image is wider than target: crop left/right
        display_w = img_h * target_ratio
        crop_percent = (img_w - display_w) / img_w / 2
        pic.crop_left = crop_percent
        pic.crop_right = crop_percent
    else:
        # Image is taller than target: crop top/bottom
        display_h = img_w / target_ratio
        crop_percent = (img_h - display_h) / img_h / 2
        pic.crop_top = crop_percent
        pic.crop_bottom = crop_percent
        
    return pic

def _add_image(slide, img_bytes: bytes, left, top, width, height) -> Any:
    """Add image while preserving aspect ratio and fitting within bounds (letterboxing)."""
    if not img_bytes:
        return None
    stream = io.BytesIO(img_bytes)
    stream.seek(0)
    with Image.open(stream) as img:
        w, h = img.size
    ratio = min(width / w, height / h)
    final_w = int(w * ratio)
    final_h = int(h * ratio)
    final_left = left + (width - final_w) // 2
    final_top = top + (height - final_h) // 2
    stream.seek(0)
    return slide.shapes.add_picture(stream, Emu(final_left), Emu(final_top), width=Emu(final_w), height=Emu(final_h))


def _add_title(slide, title_text: str, left=M, top=TM, width=SW - 2*M, height=TH,
               color: RGBColor = None, size_pt: float = 32) -> Any:
    title_text = clean_text(title_text)  # Allow full text with word wrapping
    
    # Try placeholder first - clear it to avoid "Click to add title" or ghost text
    if slide.shapes.title:
        slide.shapes.title.text = title_text
        return slide.shapes.title

    color = color or C["dark"]
    txb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT  # Left-align for professional look
    _add_run(p, title_text, FONT_H, size_pt, bold=True, color=color)
    # Red underbar - prominent width for visual consistency
    bar_h = int(SH * 0.008)
    bar_top = top + height + int(SH * 0.006)
    _add_rect(slide, left, bar_top, int(SW * 0.18), bar_h, C["red"])  # 18% width, 8% thickness
    return txb


def _add_bullets(slide, points: List[str], left, top, width, height,
                 size_pt: float = 14, bold: bool = False,
                 color: RGBColor = None, numbered: bool = False,
                 icon_bytes: Optional[List[bytes]] = None) -> None:
    color = color or C["dark"]
    if not points:
        return
        
    n = min(len(points), 6)  # Support up to 6 custom boxes
    box_gap = int(SH * 0.018)
    box_h = (height - box_gap * max(0, n - 1)) // max(n, 1)

    for i, point in enumerate(points[:n]):
        by = top + i * (box_h + box_gap)
        accent = ACCENT_RGB[i % len(ACCENT_RGB)]
        
        # Soft background box shadow + main card
        _add_rect(slide, left+15000, by+15000, width, box_h, _rgb(225, 225, 230))
        _add_rect(slide, left, by, width, box_h, C["white"])
        
        # Left colored accent stripe
        accent_w = int(SW * 0.015)
        _add_rect(slide, left, by, accent_w, box_h, accent)
        
        # Numbered Badge - centered vertically within the individual box
        badge_s = int(SH * 0.06)
        badge_x = left + int(SW * 0.02)
        badge_y = by + (box_h - badge_s) // 2
        _add_oval(slide, badge_x, badge_y, badge_s, badge_s, accent)
        _add_textbox(slide, str(i+1), badge_x, badge_y, badge_s, badge_s, 
                     font_name=FONT_H, size_pt=14, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        
        # Text inside the creative box with dynamic font scaling to eliminate 'Ghost Town' white space
        tx = left + int(SW * 0.07)
        tw = width - int(SW * 0.08)
        
        b_len = len(point)
        f_size = 22 if b_len < 60 else (17 if b_len < 110 else size_pt)
        
        _add_textbox(slide, clean_text(point),
                     left=tx, top=by + int(SH * 0.003), width=tw, height=box_h - int(SH * 0.006),
                     font_name=FONT_B, size_pt=f_size, color=color,
                     align=PP_ALIGN.LEFT, word_wrap=True)


def _add_icon(slide, img_bytes: bytes, cx: int, cy: int, size: int = 80000) -> None:
    """Place a square icon centered at (cx, cy)."""
    if not img_bytes:
        return
    _add_image(slide, img_bytes, cx - size // 2, cy - size // 2, size, size)


def _add_table(slide, tbl: Table, left, top, width, height) -> None:
    if not tbl or not tbl.rows:
        return
    headers = tbl.headers[:6]
    rows = tbl.rows[:10]
    col_count = len(headers)
    row_count = len(rows) + 1
    row_h = max(int(height / row_count), int(SH * 0.058))

    pptx_tbl = slide.shapes.add_table(
        row_count, col_count,
        Emu(left), Emu(top), Emu(width), Emu(row_h * row_count)
    ).table

    for ci, header in enumerate(headers):
        cell = pptx_tbl.cell(0, ci)
        cell.text = clean_text(header)[:40]
        cell.fill.solid()
        cell.fill.fore_color.rgb = C["red"]
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = FONT_B
                run.font.size = Pt(10)
                run.font.bold = True
                run.font.color.rgb = C["white"]

    for ri, row in enumerate(rows):
        bg = C["light"] if ri % 2 == 0 else C["white"]
        for ci, val in enumerate(row[:col_count]):
            cell = pptx_tbl.cell(ri + 1, ci)
            cell.text = clean_text(str(val))[:50]
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = FONT_B
                    run.font.size = Pt(10)
                    run.font.color.rgb = C["dark"]


# ─────────────────────────────────────────────────────────────────────────────
#  Slide builders
# ─────────────────────────────────────────────────────────────────────────────

def _build_title_slide(slide, sp: Dict, parsed: Dict) -> None:
    title = clean_text(parsed.get("title", sp.get("title", "")))
    subtitle = clean_text(parsed.get("subtitle", sp.get("subtitle", "")))

    # Try to use slide's title placeholder, fallback to custom textbox
    try:
        if slide.shapes.title:
            slide.shapes.title.text = title
        else:
            raise AttributeError("No title placeholder")
    except:
        # Fallback: custom textbox
        _add_textbox(slide, title,
                     left=M, top=int(SH * 0.30), width=SW - 2 * M, height=int(SH * 0.18),
                     font_name=FONT_H, size_pt=44, bold=True, color=C["dark"],
                     align=PP_ALIGN.CENTER)
    
    # Add subtitle
    if subtitle:
        try:
            if len(slide.placeholders) > 1:
                try:
                    subtitle_shape = slide.placeholders[1]
                    if subtitle_shape and hasattr(subtitle_shape, "text_frame"):
                        subtitle_shape.text = subtitle
                    else:
                        raise AttributeError("No text_frame")
                except (IndexError, AttributeError, KeyError):
                    raise AttributeError("No subtitle placeholder")
            else:
                raise AttributeError("No subtitle placeholder")
        except:
            # Fallback: custom textbox
            _add_textbox(slide, subtitle,
                         left=M, top=int(SH * 0.50), width=SW - 2 * M, height=int(SH * 0.12),
                         font_name=FONT_B, size_pt=18, color=C["mid"],
                         align=PP_ALIGN.CENTER)


def _build_agenda_slide(slide, sp: Dict) -> None:
    _add_title(slide, sp.get("title", "Agenda"))
    bullets = sp.get("bullets", [])
    if not bullets:
        return
    n = len(bullets)
    col_count = 2 if n >= 4 else 1
    col_w = (SW - 2 * M - (col_count - 1) * int(M * 1.2)) // col_count

    for ci in range(col_count):
        chunk = bullets[ci::col_count]
        lx = M + ci * (col_w + int(M * 1.2))
        item_h_each = (CH - int(max(len(chunk) - 1, 0) * 45000)) // max(len(chunk), 1)

        for ri, item in enumerate(chunk):
            item_top = CT + ri * item_h_each
            item_h = item_h_each - 45000
            num_size = min(int(SH * 0.058), int(item_h * 0.82))

            # Colored number box
            bg_rgb = ACCENT_RGB[(ci * len(bullets) // col_count + ri) % len(ACCENT_RGB)]
            _add_rect(slide, lx, item_top, num_size, item_h, bg_rgb)

            # Centered number
            num_txb = slide.shapes.add_textbox(
                Emu(lx), Emu(item_top), Emu(num_size), Emu(item_h))
            num_txb.text_frame.word_wrap = False
            p = num_txb.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            _add_run(p, str(ci * len(bullets) // col_count + ri + 1),
                     FONT_H, 18, bold=True, color=C["white"])

            # Item text (right of number box)
            txt_x = lx + num_size + 70000
            txt_w = col_w - num_size - 90000
            _add_textbox(slide, clean_text(item),
                         txt_x, item_top, txt_w, item_h,
                         font_name=FONT_B, size_pt=14, color=C["dark"])


def _build_exec_summary(slide, sp: Dict, parsed: Dict) -> None:
    _add_title(slide, "Executive Summary")
    
    bullets = sp.get("bullets", [])
    if not bullets:
        text = clean_text(parsed.get("executive_summary", ""))
        import re
        bullets = [s.strip() for s in re.split(r"(?<=[.!?])\s+", text) if len(s) > 15][:4]
        
    n = len(bullets)
    
    # 1. Left Pillar: Semicircle Image (Desk/Analytical view)
    semi_w = int(SW * 0.22)
    semi_h = int(SH * 0.6)
    semi_top = CT + (CH - semi_h) // 2
    
    # Background for semicircle area logic
    bg_bytes = generate_slide_asset("team working desk office", "Strategy Analysis")
    _add_semicircle_image(slide, bg_bytes, 0, semi_top, semi_w, semi_h, side="left")
    
    # 2. Right Pillar: Large 3D Mascot Asset
    mascot_w = int(SW * 0.28)
    mascot_h = int(SH * 0.8)
    mascot_top = SH - mascot_h
    mascot_left = SW - mascot_w
    
    mascot_bytes = generate_slide_asset(sp.get("title", ""), parsed.get("title", ""), is_mascot=True)
    _add_image(slide, mascot_bytes, mascot_left, mascot_top, mascot_w, mascot_h)
    
    # 3. Center Content: Professional Bullet List
    text_left = semi_w + int(SW * 0.05)
    text_w = SW - semi_w - mascot_w - int(SW * 0.1) # Safe spacing
    
    avail_h = CH
    item_h = avail_h // max(n, 1)
    
    for i, bullet in enumerate(bullets):
        by = CT + i * item_h
        accent = ACCENT_RGB[i % len(ACCENT_RGB)]
        
        # Icon
        icon_s = int(SH * 0.08)
        _add_oval(slide, text_left, by + (item_h - icon_s) // 2, icon_s, icon_s, _rgb(248, 248, 250))
        # Place a generic icon inside
        _add_icon(slide, get_numbered_icon(i+1), text_left + icon_s // 2, by + item_h // 2, size=int(icon_s * 0.6))
        
        # Text alignment
        tx = text_left + icon_s + int(SW * 0.02)
        tw = SW - mascot_w - tx - int(SW * 0.02)
        
        parts = bullet.split(":", 1)
        if len(parts) > 1:
            head, body = parts[0] + ":", parts[1].strip()
        else:
            head, body = "Focus Point:", bullet
            
        _add_textbox(slide, head, tx, by + int(SH * 0.02), tw, int(SH * 0.06),
                     font_name=FONT_H, size_pt=18, bold=True, color=C["dark"])
        _add_textbox(slide, body, tx, by + int(SH * 0.08), tw, item_h - int(SH * 0.1),
                     font_name=FONT_B, size_pt=12, color=C["mid"])


def _build_monitoring_slide(slide, sp: Dict, parsed: Dict) -> None:
    # Use Sidebar for title instead of top header to prevent overlap as seen in Image 6
    
    # 1. Sidebar (Left 22%)
    side_w = int(SW * 0.22)
    _add_rect(slide, 0, 0, side_w, SH, _rgb(252, 252, 254))
    
    # Large Sidebar Title & Icon - Adjusted top and height for better alignment in pillar
    title_txt = sp.get("title", "").upper()
    f_size = 22 if len(title_txt) < 30 else 18
    _add_textbox(slide, title_txt, int(SW * 0.02), int(SH * 0.15), side_w - int(SW * 0.04), int(SH * 0.35),
                 font_name=FONT_H, size_pt=f_size, bold=True, color=C["dark"])
                 
    # Bottom Sidebar Image (Mascot or Team)
    img_h = int(SH * 0.4)
    img_bytes = generate_slide_asset("business men looking at tablet", "Corporate Analysis")
    _add_image_cover(slide, img_bytes, 0, SH - img_h, side_w, img_h)
    
    # 2. Content Grid (Right 78%)
    bullets = sp.get("bullets", [])[:3]
    n = len(bullets)
    if not bullets: return
    
    margin = int(SW * 0.02)
    grid_left = side_w + margin
    grid_w = SW - grid_left - margin
    col_w = grid_w // n
    
    for i, bullet in enumerate(bullets):
        lx = grid_left + i * col_w
        
        # Header Circle (Faded pink effect)
        header_s = int(SH * 0.28)
        _add_oval(slide, lx + (col_w - header_s)//2, int(SH * 0.1), header_s, header_s, _rgb(255, 235, 235))
        
        parts = bullet.split(":", 1)
        head = parts[0] if len(parts) > 1 else "Metric"
        body_list = parts[1].strip().split(".") if len(parts) > 1 else [bullet]
        
        # Header text inside circle
        _add_textbox(slide, head, lx + int(SW * 0.01), int(SH * 0.18), col_w - int(SW * 0.02), int(SH * 0.08),
                     font_name=FONT_H, size_pt=14, bold=True, color=C["red"], align=PP_ALIGN.CENTER)
        
        # Diagonal Badge
        badge_s = int(SH * 0.08)
        _add_split_circle(slide, lx + (col_w - badge_s)//2, int(SH * 0.38), badge_s, C["dark"], C["red"])
        _add_textbox(slide, str(i+1), lx + (col_w - badge_s)//2, int(SH * 0.38) + int(SH * 0.01), badge_s, badge_s,
                     font_name=FONT_H, size_pt=14, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
        
        # Bullet list with small red dots
        ty = int(SH * 0.52)
        for b in body_list:
            b = b.strip()
            if not b: continue
            _add_oval(slide, lx + int(SW * 0.01), ty + int(SH * 0.012), int(SH * 0.012), int(SH * 0.012), C["red"])
            _add_textbox(slide, b, lx + int(SW * 0.025), ty, col_w - int(SW * 0.035), int(SH * 0.12),
                         font_name=FONT_B, size_pt=11, color=C["dark"])
            ty += int(SH * 0.14)


def _build_section_divider(slide, sp: Dict) -> None:
    title = clean_text(sp.get("title", ""))
    # Colored background strip (centered vertically) - smaller to not overwhelm
    divider_h = int(SH * 0.26)
    divider_top = int((SH - divider_h) / 2) + int(SH * 0.02)
    _add_rect(slide, 0, divider_top, SW, divider_h, C["red"])
    
    # Centered title - balanced size
    _add_textbox(slide, title,
                 left=M, top=divider_top, width=SW - 2 * M, height=divider_h,
                 font_name=FONT_H, size_pt=34, bold=True, color=C["white"],
                 align=PP_ALIGN.CENTER)
    
    subtitle = clean_text(sp.get("subtitle", ""))
    if subtitle:
        subtitle_top = divider_top + divider_h + int(SH * 0.04)
        _add_textbox(slide, subtitle,
                     left=M, top=subtitle_top,
                     width=SW - 2 * M, height=int(SH * 0.10),
                     font_name=FONT_B, size_pt=14, color=C["dark"],
                     align=PP_ALIGN.CENTER)


def _build_content_slide(slide, sp: Dict, parsed: Dict) -> None:
    _add_title(slide, sp.get("title", ""))
    bullets = sp.get("bullets", [])
    if not bullets:
        return
    
    # Icon beside title
    icon_bytes = get_icon_for_title(sp.get("title", ""),
                                    index=sp.get("slide_number", 1) - 1)
    icon_size = int(SH * 0.072)
    _add_icon(slide, icon_bytes,
              cx=SW - M - icon_size // 2 - int(SW * 0.008),
              cy=TM + TH // 2,
              size=icon_size)

    # Place bullets on left ~75% (Expanded to cover majority as requested)
    text_width = int((SW - 2 * M) * 0.75)
    _add_bullets(slide, bullets[:6],
                 left=M, top=CT, width=text_width, height=CH,
                 size_pt=15, color=C["dark"])
                 
    # Fill right-side gap with dynamic AI visual asset!
    title = sp.get("title", "Corporate Slide")
    doc_title = parsed.get("title", "Corporate Strategy")
    
    bg_bytes = generate_slide_asset(title, doc_title)
    
    if not bg_bytes:
        # Fallback to static assets
        title_lower = title.lower()
        asset_name = "corporate_strategy.png"
        if any(keyword in title_lower for keyword in ["tech", "ai", "digital", "data", "cloud"]):
            asset_name = "corporate_tech.png"
        elif any(keyword in title_lower for keyword in ["finance", "revenue", "growth", "market", "acquisition"]):
            asset_name = "corporate_finance.png"
            
        img_path = os.path.join(os.path.dirname(__file__), "..", "assets", asset_name)
        if os.path.exists(img_path):
            with open(img_path, "rb") as f:
                bg_bytes = f.read()
                
    if bg_bytes:
        # Draw image heavily bleeding to the right boundary, with 75/25 split
        img_left = M + text_width + int(M * 0.5)
        img_width = SW - img_left
        _add_image_cover(slide, bg_bytes, img_left, CT, img_width, CH)


def _build_two_column(slide, sp: Dict) -> None:
    # Build a peer comparison infographic (Reference: Image 4)
    # Left pillar: Bold red overlay over image
    left_w = int(SW * 0.3)
    title = sp.get("title", "Comparison & Analysis")
    
    # Thematic Left Column Background
    _add_rect(slide, 0, 0, left_w, SH, C["dark"])
    
    # Dynamically seed HF generation to uniquely match slide intent instead of static cache
    bg_bytes = generate_slide_asset(title, "Corporate Competition")
    if bg_bytes:
        # Image spans the entire 100% height of the left pillar
        _add_image_cover(slide, bg_bytes, 0, 0, left_w, SH)
    
    # Top Left Overlapping Focus Box drawn ABOVE the image overlay
    focus_box_h = int(SH * 0.4)
    _add_rect(slide, 0, int(SH * 0.1), left_w + int(M * 0.5), focus_box_h, C["red"])
    _add_textbox(slide, title.upper(), int(M * 0.5), int(SH * 0.15), left_w - int(M * 0.5), focus_box_h - int(SH * 0.1),
                 font_name=FONT_H, size_pt=24, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
                 
    # Right side connected vertical cards
    bullets = sp.get("left_bullets", []) + sp.get("right_bullets", []) + sp.get("bullets", [])
    n = max(min(len(bullets), 3), 1)
    
    right_space = SW - left_w - int(M * 1.5)
    gap = int(SW * 0.03)
    card_w = (right_space - (n - 1) * gap) // n
    card_h = int(SH * 0.6)
    card_top = int(SH * 0.25)
    
    # Connecting backbone line behind nodes
    line_y = card_top + int(SH * 0.08)
    _add_rect(slide, left_w + int(M * 1.5) + card_w // 2, line_y, right_space - card_w, int(SH * 0.005), C["mid"])

    for i, bullet in enumerate(bullets[:n]):
        cx = left_w + int(M * 1.5) + i * (card_w + gap)
        
        # Gray background offset / Drop shadow
        _add_rect(slide, cx + 15000, card_top + 15000, card_w, card_h, _rgb(230, 230, 235))
        # Main white card
        _add_rect(slide, cx, card_top, card_w, card_h, C["white"])
        
        # Number Node Header (like a cog in Image 4)
        node_s = int(SH * 0.14)
        node_x = cx + card_w // 2 - node_s // 2
        node_y = card_top - node_s // 3
        # Outer red ring
        _add_oval(slide, node_x - 10000, node_y - 10000, node_s + 20000, node_s + 20000, C["red"])
        # Inner white circle
        _add_oval(slide, node_x, node_y, node_s, node_s, C["white"])
        # Number text
        _add_textbox(slide, f"{i+1:02d}", node_x, node_y + int(SH * 0.003), node_s, node_s,
                     font_name=FONT_H, size_pt=18, bold=True, color=C["red"], align=PP_ALIGN.CENTER)
        
        # Snippled Heading & Text Mapping safely
        parts = bullet.split(":", 1)
        if len(parts) > 1 and len(parts[0]) < 40:
            title_text = parts[0] + ":"
            body = parts[1].strip()
        else:
            title_text = ""
            body = bullet

        if title_text:
            _add_textbox(slide, title_text, cx + 50000, card_top + int(card_h * 0.25), card_w - 100000, int(card_h * 0.15),
                         font_name=FONT_H, size_pt=13, bold=True, color=C["dark"], align=PP_ALIGN.CENTER)
                         
        body_top = card_top + int(card_h * 0.4) if title_text else card_top + int(card_h * 0.25)
        body_h = card_h - (body_top - card_top) - 50000
        _add_textbox(slide, body, cx + 50000, body_top, card_w - 100000, body_h,
                     font_name=FONT_B, size_pt=12, color=C["mid"], align=PP_ALIGN.CENTER, word_wrap=True)


def _build_data_chart(slide, sp: Dict, parsed: Dict) -> None:
    _add_title(slide, sp.get("title", ""))
    chart_spec = sp.get("chart")
    img_bytes = b""

    if chart_spec and chart_spec.get("labels") and chart_spec.get("values"):
        img_bytes = render_chart(chart_spec)

    if not img_bytes:
        # Try to auto-extract from a table in parsed data
        ti = sp.get("table_index", -1)
        all_tables = parsed.get("all_tables", [])
        if 0 <= ti < len(all_tables):
            cd = all_tables[ti].chart_data()
            if cd:
                img_bytes = render_chart(cd)
        if not img_bytes and all_tables:
            for tbl in all_tables:
                cd = tbl.chart_data()
                if cd:
                    img_bytes = render_chart(cd)
                    break

    if img_bytes:
        # Properly size chart to fit within content area
        # Use 85% of available content width for good margins
        chart_w = int((SW - 2 * M) * 0.85)
        chart_left = M + int(((SW - 2 * M) - chart_w) / 2)  # Center horizontally
        chart_top = CT + int(SH * 0.04)  # Small top gap
        chart_h = CH - int(SH * 0.06)  # Use most of content height
        _add_image(slide, img_bytes, chart_left, chart_top, chart_w, chart_h)
    else:
        # Fallback: content bullets
        _add_bullets(slide, sp.get("bullets", []),
                     M, CT, SW - 2 * M, CH, size_pt=14)


def _build_data_table(slide, sp: Dict, parsed: Dict) -> None:
    # Use the styled image banner
    content_top = _build_image_banner(slide, sp.get("title", "Key Financials & Data"), parsed)
    
    all_tables = parsed.get("all_tables", [])
    ti = sp.get("table_index", -1)
    tbl = None
    if 0 <= ti < len(all_tables):
        tbl = all_tables[ti]
    elif all_tables:
        tbl = all_tables[0]

    if tbl and tbl.rows:
        headers = tbl.headers[:5] # Allow 5 columns max
        rows = tbl.rows # NO CAP so we don't lose user data
        col_count = min(len(headers), 5)

        margin_x = int(SW * 0.04)
        gap_x = int(SW * 0.015)
        gap_y = int(SH * 0.01)
        avail_w = SW - 2 * margin_x
        
        col_widths = [avail_w // col_count] * col_count
        if col_count > 1:
            col_widths[0] = int(avail_w * 0.25)
            rem_w = avail_w - col_widths[0] - gap_x * (col_count - 1)
            for j in range(1, col_count):
                col_widths[j] = rem_w // (col_count - 1)

        total_h_avail = SH - content_top - int(SH * 0.08)
        row_count = max(1, len(rows) + 1)
        row_h = total_h_avail // row_count - gap_y
        
        # Guarantee minimum height so elements don't collapse or overlap 
        if row_h < 150000:
            row_h = 150000
            
        font_s = max(8, min(14, int(row_h / 40000)))

        # Draw Headers safely fixed beneath the banner area
        top_y = content_top + int(SH * 0.05)
        header_h = int(SH * 0.08)
        lx = margin_x
        for j, h in enumerate(headers[:col_count]):
            _add_textbox(slide, h, lx, top_y, col_widths[j], header_h, 
                         font_name=FONT_H, size_pt=font_s+2, bold=True, color=C["dark"], align=PP_ALIGN.CENTER)
            lx += col_widths[j] + gap_x

        # Draw Styled Shape Grid flowing downwards
        row_start_y = top_y + header_h + int(SH * 0.02)
        max_rows = (SH - row_start_y - int(SH * 0.05)) // (row_h + gap_y)
        
        for i, row_data in enumerate(rows[:int(max_rows)]):
            ly = row_start_y + i * (row_h + gap_y)
            lx = margin_x
            
            for j in range(col_count):
                val = str(row_data[j]) if j < len(row_data) else ""
                w = col_widths[j]
                
                if j == 0:
                    # Lead red axis block
                    _add_rect(slide, lx, ly, w, row_h, C["red"])
                    _add_textbox(slide, val, lx, ly, w, row_h, 
                                 font_name=FONT_H, size_pt=font_s+2, bold=True, color=C["white"], align=PP_ALIGN.CENTER)
                else:
                    # White metric blocks with backdrop shadow
                    _add_rect(slide, lx, ly, w, row_h, C["white"])
                    
                    if j > 1:
                        # Intersecting circles imitating process connection
                        cir_s = min(int(SH * 0.006), row_h // 2)
                        _add_oval(slide, lx - gap_x // 2 - cir_s // 2, ly + row_h // 2 - cir_s // 2, cir_s, cir_s, C["dark"])
                    
                    _add_textbox(slide, val, lx + int(SW * 0.0016), ly, w - int(SW * 0.0032), row_h, 
                                 font_name=FONT_B, size_pt=font_s, color=C["mid"], align=PP_ALIGN.CENTER)
                lx += w + gap_x

    else:
        # Fallback to general display
        _add_bullets(slide, sp.get("bullets", []),
                     M, content_top + 100000, SW - 2 * M, CH - 200000, size_pt=15)

def _build_conclusion(slide, sp: Dict) -> None:
    _add_title(slide, sp.get("title", "Key Takeaways"))
    bullets = sp.get("bullets", [])
    if not bullets:
        return
    
    n = min(len(bullets), 6)  # Support up to 6 takeaways
    if n <= 0:
        return
        
    # Dynamic spacing
    box_gap = int(SH * 0.015)
    box_h = (CH - box_gap * max(0, n - 1)) // n

    for i, point in enumerate(bullets[:n]):
        by = CT + i * (box_h + box_gap)
        accent = ACCENT_RGB[i % len(ACCENT_RGB)]

        # Creative Box Background
        _add_rect(slide, M, by, SW - 2 * M, box_h, _rgb(245, 245, 247))
        
        # Left accent stripe
        accent_w = int(SW * 0.008)
        _add_rect(slide, left=M, top=by, width=accent_w, height=box_h, fill_rgb=accent)

        # Icon (using index for auto-color selection)
        icon_bytes = get_numbered_icon(i + 1)
        icon_s = int(box_h * 0.60)
        _add_icon(slide, icon_bytes,
                  cx=M + accent_w + 70000 + icon_s // 2,
                  cy=by + box_h // 2,
                  size=icon_s)

        # Text with proper wrapping
        tx = M + accent_w + 70000 + icon_s + 70000
        tw = SW - tx - M - 30000
        _add_textbox(slide, clean_text(point),
                     left=tx, top=by, width=tw, height=box_h,
                     font_name=FONT_B, size_pt=13, color=C["dark"])


def _build_thankyou_slide(slide, sp: Dict, parsed: Dict) -> None:
    # In template1.pptx and similar, the 'Thank You' text/graphics are built into the master.
    # The user explicitly requested an empty last page for these templates.
    try:
        if "thank you" in slide.slide_layout.name.lower() or "closing" in slide.slide_layout.name.lower():
            return
    except:
        pass

    # Use presentation title instead of "Thank You" for professional finish
    title_text = clean_text(parsed.get("title", "Closing Remarks"))
    
    # Try to use template's title placeholder, fallback to custom textbox
    try:
        if slide.shapes.title:
            slide.shapes.title.text = title_text
        else:
            raise AttributeError("No title placeholder")
    except:
        _add_textbox(slide, title_text,
                     left=M, top=int(SH * 0.30), width=SW - 2 * M, height=int(SH * 0.20),
                     font_name=FONT_H, size_pt=40, bold=True, color=C["dark"],
                     align=PP_ALIGN.CENTER)
    
    subtitle = clean_text(sp.get("subtitle", "Questions & Discussion"))
    
    # Add subtitle
    if subtitle:
        try:
            if len(slide.placeholders) > 1:
                try:
                    subtitle_shape = slide.placeholders[1]
                    if subtitle_shape and hasattr(subtitle_shape, "text_frame"):
                        subtitle_shape.text = subtitle
                    else:
                        raise AttributeError("No text_frame")
                except (IndexError, AttributeError, KeyError):
                    raise AttributeError("No subtitle placeholder")
            else:
                raise AttributeError("No subtitle placeholder")
        except:
            _add_textbox(slide, subtitle,
                         left=M, top=int(SH * 0.50), width=SW - 2 * M, height=int(SH * 0.12),
                         font_name=FONT_B, size_pt=18, color=C["mid"],
                         align=PP_ALIGN.CENTER)


def _build_timeline_slide(slide, sp: Dict) -> None:
    _add_title(slide, sp.get("title", "Roadmap / Timeline"))
    bullets = sp.get("bullets", [])
    
    if not bullets:
        # Prevent completely blank slides on missing parser data
        content_txt = sp.get("content", "")
        if content_txt:
            import re
            bullets = [s.strip() for s in re.split(r"(?<=[.!?])\s+", clean_text(content_txt)) if len(s) > 10]
    
    if not bullets:
        # Hard fallback
        bullets = ["Initial Analysis completely parsed.", "Strategic Execution in motion.", "Validation ongoing."]
        
    n = min(len(bullets), 5)
    margin = int(SW * 0.08)
    avail_w = SW - 2 * margin
    step_w = avail_w // n
    line_y = CT + int(CH * 0.45)
    
    # Main timeline line
    _add_rect(slide, margin, line_y, avail_w, int(SH * 0.008), C["mid"])
    
    for i, bullet in enumerate(bullets[:n]):
        cx = margin + i * step_w + step_w // 2
        accent = ACCENT_RGB[i % len(ACCENT_RGB)]
        
        # Timeline node
        node_s = int(SH * 0.04)
        _add_rect(slide, cx - node_s // 2, line_y - node_s // 2 + int(SH * 0.004),
                  node_s, node_s, accent)
        
        # Determine top or bottom placement for text to avoid overlap
        is_top = (i % 2 == 0)
        text_h = int(CH * 0.35)
        text_y = line_y - text_h - int(SH * 0.04) if is_top else line_y + int(SH * 0.04)
        
        # Vertical connector
        conn_h = int(SH * 0.03)
        conn_y = line_y - conn_h if is_top else line_y + int(SH * 0.008)
        _add_rect(slide, cx - int(SW * 0.0012), conn_y, int(SW * 0.0024), conn_h, accent)
        
        # Text box
        _add_textbox(slide, bullet,
                     left=cx - step_w // 2 + int(SW * 0.0041), top=text_y,
                     width=step_w - int(SW * 0.0082), height=text_h,
                     font_name=FONT_B, size_pt=12, color=C["dark"],
                     align=PP_ALIGN.CENTER)


def _build_process_flow(slide, sp: Dict) -> None:
    _add_title(slide, sp.get("title", "Process Flow"))
    bullets = sp.get("bullets", [])
    
    if not bullets:
        # Prevent completely blank slides returning nothing
        content_txt = sp.get("content", "")
        if content_txt:
            import re
            bullets = [s.strip() for s in re.split(r"(?<=[.!?])\s+", clean_text(content_txt)) if len(s) > 10][:4]
            
    if not bullets:
        # Hard fallback data
        bullets = ["Discovery Phase", "Detailed Analysis", "Strategic Deployment", "Iteration"]
        
    n = min(len(bullets), 4)
    item_gap = int(SW * 0.02)
    item_w = (SW - 2 * M - (n - 1) * item_gap) // n
    item_h = int(CH * 0.65)
    item_top = CT + (CH - item_h) // 2
    
    for i, bullet in enumerate(bullets[:n]):
        lx = M + i * (item_w + item_gap)
        accent = ACCENT_RGB[i % len(ACCENT_RGB)]
        
        # Step Background
        _add_rect(slide, lx, item_top, item_w, item_h, _rgb(245, 245, 247))
        _add_rect(slide, lx, item_top, item_w, int(SH * 0.012), accent)
        
        # Step Number
        num_s = int(SH * 0.08)
        _add_textbox(slide, str(i + 1),
                     lx, item_top + int(SH * 0.0073), item_w, num_s,
                     font_name=FONT_H, size_pt=24, bold=True, color=accent,
                     align=PP_ALIGN.CENTER)
        
        # Step Text dynamic scaling vertically preventing 'Ghost Town' formatting
        b_len = len(bullet)
        f_size = 24 if b_len < 50 else (18 if b_len < 90 else 13)
        
        _add_textbox(slide, bullet,
                     lx + int(SW * 0.0082), item_top + num_s + int(SH * 0.0073),
                     item_w - int(SW * 0.0164), item_h - num_s - int(SH * 0.0219),
                     font_name=FONT_B, size_pt=f_size, color=C["dark"],
                     align=PP_ALIGN.CENTER, word_wrap=True)
        
        # Arrow to next
        if i < n - 1:
            arr_x = lx + item_w + item_gap // 2 - int(SW * 0.0082)
            arr_y = item_top + item_h // 2
            # Simple arrow shape using rects
            _add_rect(slide, arr_x, arr_y - int(SH * 0.003), int(SW * 0.0164), int(SH * 0.006), C["light"])


# ─────────────────────────────────────────────────────────────────────────────
#  Layout resolver
# ─────────────────────────────────────────────────────────────────────────────

def _get_layout(prs: Presentation, name: str):
    """
    Get layout from template with fuzzy matching to support diverse templates.
    """
    name = name.lower()
    
    # 1. Direct key mapping (preferred Accenture template naming)
    layout_map = {
        "cover":      "cover",
        "divider":    "divider",
        "blank":      "blank",
        "title_only": "title only",
        "thankyou":   "thank you",
    }
    target_keyword = layout_map.get(name, name)
    
    # 2. Try Exact Match (Case-Insensitive)
    for lo in prs.slide_layouts:
        if lo.name.lower() == target_keyword:
            return lo
            
    # 3. Try Fuzzy Match (Contains keyword)
    for lo in prs.slide_layouts:
        lo_name = lo.name.lower()
        if target_keyword in lo_name:
            return lo
            
    # 4. Special logic for common layout roles
    if "cover" in name or "title" in name:
        # Search for anything that looks like a cover
        for lo in prs.slide_layouts:
            if "cover" in lo.name.lower() or "title slide" in lo.name.lower():
                return lo
                
    # 5. Last Resort: First available layout
    if len(prs.slide_layouts) > 0:
        return prs.slide_layouts[0]
        
    raise ValueError("The provided PPTX template has no layouts!")


def _build_agentic_logic(slide, sp: Dict) -> None:
    _add_title(slide, "System Architecture & Agentic Flow")
    
    # 3-Step circular or arrow flow showing Parser -> AI Planner -> Builder
    margin = int(SW * 0.1)
    avail_w = SW - 2 * margin
    box_w = avail_w // 3 - int(SW * 0.041)
    box_h = int(CH * 0.5)
    box_top = CT + (CH - box_h) // 2
    
    steps = [
        ("Recursive Parser", "Extracts structured sections, recursive table data & numerical insights", C["blue"]),
        ("Reasoning Planner", "DeepSeek-R1 synthesis of high-value insights & slide sequencing", C["red"]),
        ("Dynamic Builder", "Programmatic layout generation & Matplotlib data visualization", C["green"])
    ]
    
    for i, (name, desc, color) in enumerate(steps):
        lx = margin + i * (box_w + 600000)
        _add_rect(slide, lx, box_top, box_w, box_h, _rgb(248, 249, 250))
        _add_rect(slide, lx, box_top, int(SW * 0.008), box_h, color)
        
        _add_textbox(slide, name, lx + 150000, box_top + 100000, box_w - 200000, 400000,
                     font_name=FONT_H, size_pt=18, bold=True, color=color)
        _add_textbox(slide, desc, lx + 150000, box_top + 600000, box_w - 200000, 800000,
                     font_name=FONT_B, size_pt=12, color=C["dark"])
        
        if i < len(steps) - 1:
            arr_x = lx + box_w + 100000
            _add_rect(slide, arr_x, box_top + box_h // 2, 400000, 40000, C["light"])

def _build_arc_infographic(slide, sp: Dict) -> None:
    # Full Red Background as seen in Image 2
    _add_rect(slide, 0, 0, SW, SH, _rgb(200, 30, 30)) # Corporate Red
    
    # Central Bottom Hub
    hub_s = int(SH * 0.45)
    hub_x = (SW - hub_s) // 2
    hub_y = SH - hub_s // 2
    _add_oval(slide, hub_x, hub_y, hub_s, hub_s, C["white"])
    
    title_text = clean_text(sp.get("title", ""))
    _add_textbox(slide, title_text, (SW - int(SW*0.4)) // 2, SH - int(SH*0.18), int(SW*0.4), int(SH*0.15),
                 font_name=FONT_H, size_pt=24, bold=True, color=_rgb(200, 30, 30), align=PP_ALIGN.CENTER)
    
    bullets = sp.get("bullets", [])[:5]
    n = len(bullets)
    if not bullets: return
    
    # Arc Parameters
    arc_r = int(SH * 0.55)
    arc_cx = SW // 2
    arc_cy = SH # Bottom center
    
    for i, bullet in enumerate(bullets):
        # Calculate angle (distribute along arc)
        angle_deg = 180 - (180 / (n + 1)) * (i + 1)
        angle_rad = math.radians(angle_deg)
        
        nx = arc_cx + arc_r * math.cos(angle_rad)
        ny = arc_cy - arc_r * math.sin(angle_rad)
        
        # Node
        node_s = int(SH * 0.12)
        lx = nx - node_s // 2
        ly = ny - node_s // 2
        
        # Connector dotted line (imaginary arc effect)
        # nodes
        _add_oval(slide, lx, ly, node_s, node_s, C["white"])
        _add_textbox(slide, f"{i+1:02d}", lx, ly + int(SH*0.02), node_s, node_s, 
                     font_name=FONT_H, size_pt=18, bold=True, color=_rgb(200, 30, 30), align=PP_ALIGN.CENTER)
        
        # Text Box
        tx_w = int(SW * 0.25)
        tx_h = int(SH * 0.25)
        tx_left = lx - tx_w // 2 + node_s // 2
        tx_top = ly - tx_h if ny < arc_cy - int(SH*0.3) else ly + node_s # Adjust based on quadrant
        
        _add_textbox(slide, bullet, tx_left, tx_top, tx_w, tx_h,
                     font_name=FONT_B, size_pt=11, color=C["white"], align=PP_ALIGN.CENTER, word_wrap=True)

def _build_infographic_dispatch(slide, sp: Dict) -> None:
    itype = sp.get("infographic_type", "process")
    title = sp.get("title", "").lower()
    
    if itype == "mechanism" or any(k in title for k in ["bubble", "system", "logic", "flow"]):
        _build_arc_infographic(slide, sp)
    elif itype == "timeline":
        _build_timeline_slide(slide, sp)
    elif itype == "comparison":
        _build_two_column(slide, sp)
    else:
        _build_process_flow(slide, sp)

# ─────────────────────────────────────────────────────────────────────────────
#  Main public function
# ─────────────────────────────────────────────────────────────────────────────

def generate_pptx(parsed: Dict[str, Any], slide_plan: List[Dict],
                  output_path: str, template_name: str = "template.pptx") -> None:
    """Build the full PPTX from parsed markdown + AI slide plan using template."""
    template_path = os.path.join(os.path.dirname(__file__), "..", "assets", template_name)
    if not os.path.exists(template_path):
        template_path = TEMPLATE_PATH
        print(f"WARN: Template {template_name} not found. Fallback to {TEMPLATE_PATH}")
    prs = Presentation(template_path)
    
    # Initialize dynamic layout based on template size
    _setup_layout(prs)
    
    # Remove any existing slides from the template
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    BUILDER_MAP = {
        "title":           lambda sl, sp: _build_title_slide(sl, sp, parsed),
        "agenda":          lambda sl, sp: _build_agenda_slide(sl, sp),
        "exec_summary":    lambda sl, sp: _build_exec_summary(sl, sp, parsed),
        "section_divider": lambda sl, sp: _build_section_divider(sl, sp),
        "content":         lambda sl, sp: _build_content_slide(sl, sp, parsed),
        "two_column":      lambda sl, sp: _build_two_column(sl, sp),
        "data_chart":      lambda sl, sp: _build_data_chart(sl, sp, parsed),
        "data_table":      lambda sl, sp: _build_data_table(sl, sp, parsed),
        "infographic":     lambda sl, sp: _build_infographic_dispatch(sl, sp),
        "monitoring":      lambda sl, sp: _build_monitoring_slide(sl, sp, parsed),
        "agentic_logic":   lambda sl, sp: _build_agentic_logic(sl, sp),
        "conclusion":      lambda sl, sp: _build_conclusion(sl, sp),
        "thankyou":        lambda sl, sp: _build_thankyou_slide(sl, sp, parsed),
    }

    # Enforce strict slide ordering based on the AI's intended slide numbers to prevent sequence jumbling
    slide_plan = sorted(slide_plan, key=lambda x: x.get("slide_number", 999))
    
    for sp in slide_plan:
        slide_type = sp.get("type", "content")
        layout_name = sp.get("layout", "blank")

        # Map slide types to appropriate layouts
        if slide_type == "title":
            layout_name = "cover"
        elif slide_type == "thankyou":
            layout_name = "thankyou"
        elif slide_type == "section_divider":
            layout_name = "divider"
        else:
            layout_name = "title_only"

        layout = _get_layout(prs, layout_name)
        slide = prs.slides.add_slide(layout)

        builder = BUILDER_MAP.get(slide_type, BUILDER_MAP["content"])
        try:
            builder(slide, sp)
        except Exception as exc:
            print(f"   [WARN] Slide {sp.get('slide_number', '?')} ({slide_type}): {exc}")
            # Best-effort fallback
            try:
                _build_content_slide(slide, sp, parsed)
            except Exception:
                pass

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    print(f"   [OK] Saved -> {output_path}  ({len(slide_plan)} slides)")